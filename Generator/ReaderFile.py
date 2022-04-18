## coding: UTF-8
import pathlib
from chardet.universaldetector import UniversalDetector
import pandas as pd
from bs4 import BeautifulSoup
import docx
import pptx
from . import GenerateWordCloud
from .ReaderWeb import GenerateFromWeb

# Ref: https://kazusa-pg.com/python-detect-character-code/
def detect_character_code(file):
    detector = UniversalDetector()
    with open(file, 'rb') as f:
        detector.reset()
        for line in f.readlines():
            detector.feed(line)
            if detector.done:
                break
        detector.close()
    return detector.result['encoding']

def read_plane_text(path):
    code = detect_character_code(str(path))
    with open(path, 'r', encoding=code) as f:
        while True:
            line = f.readline()
            if line == '':
                break
            yield line.rstrip('\n')

def read_table_text(path):
    code = detect_character_code(str(path))
    path = pathlib.Path(str(path))
    if path.suffix == '.csv':
        dfs = {'csv': pd.read_csv(str(path), encoding=code)}
    elif path.suffix == '.tsv':
        dfs = {'tsv' : pd.read_table(str(path), encoding=code)}
    elif path.suffix in ['.xls', '.xlsx']:
        dfs = pd.read_excel(str(path), sheet_name=None)
    for _, df in dfs.items():
        str_cols = []
        for _, vals in df.iterrows():
            for col, val in vals.items():
                if isinstance(val, str):
                    str_cols.append(col)
            break
        if len(str_cols)>0:
            yield ' '.join(str_cols)
            df = df[str_cols]
            for line in df.values:
                yield ' '.join(line)

def read_html_text(path):
    code = detect_character_code(str(path))
    with open(path, 'r', encoding=code) as f:
        soup = BeautifulSoup(f.read(), "html.parser")
        for line in GenerateFromWeb.get_page_text_only(soup):
            yield line

def read_doc_text(path):
    doc = docx.Document(str(path))
    for para in doc.paragraphs:
        yield para.text
    for tbl in doc.tables:
        for row in tbl.rows:
            row_text = [cell.text for cell in row.cells]
            yield ' '.join(row_text)

def read_ppt_text(path):
    prs = pptx.Presentation(str(path))
    for sld in prs.slides:
        for shp in sld.shapes:
            if shp.has_text_frame:
                yield shp.text

class GenerateFromFile(GenerateWordCloud):
    def __init__(self, *args, path, **kwargs):
        super().__init__(*args, **kwargs)
        self.path = pathlib.Path(str(path))
    
    @staticmethod
    def read_text(path):
        reader = None
        # *.ppt *.pptx
        if path.suffix in ['.txt']:
            reader = read_plane_text
        elif path.suffix in ['.csv', '.tsv', '.xls', '.xlsx']:
            reader = read_table_text
        elif path.suffix in ['.htm', '.html']:
            reader = read_html_text
        elif path.suffix in ['.doc', '.docx']:
            reader = read_doc_text
        elif path.suffix in ['.ppt', '.pptx']:
            reader = read_ppt_text
        
        if reader is not None:
            for line in reader(path):
                yield line
    
    def read_text_lines(self):
        for line in self.read_text(self.path):
            yield line